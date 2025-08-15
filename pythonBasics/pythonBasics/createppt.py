from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.dml.color import RGBColor

# Create a new PowerPoint presentation
prs = Presentation()


# Helper function to add a slide with title and bullet points
def add_slide(layout, title_text, bullet_points, subtitle_text=None):
    slide = prs.slides.add_slide(layout)
    title = slide.shapes.title
    title.text = title_text
    title.text_frame.paragraphs[0].font.size = Pt(32)
    title.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)

    if subtitle_text:
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        subtitle.text_frame.paragraphs[0].font.size = Pt(18)

    content = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9), Inches(4.5))
    tf = content.text_frame
    tf.word_wrap = True
    for point in bullet_points:
        p = tf.add_paragraph()
        p.text = point
        p.font.size = Pt(16)
        p.level = 0 if not point.startswith("  ") else 1
        p.alignment = PP_ALIGN.LEFT


# Slide 1: Title Slide
title_slide_layout = prs.slide_layouts[0]
slide1 = prs.slides.add_slide(title_slide_layout)
title1 = slide1.shapes.title
title1.text = "Playwright KT Session"
title1.text_frame.paragraphs[0].font.size = Pt(44)
title1.text_frame.paragraphs[0].font.color.rgb = RGBColor(0, 51, 102)
subtitle1 = slide1.placeholders[1]
subtitle1.text = "Rapid and Reliable End-to-End Testing for Web"
subtitle1.text_frame.paragraphs[0].font.size = Pt(24)

# Slide 2: Agenda
bullet_slide_layout = prs.slide_layouts[1]
add_slide(bullet_slide_layout, "Agenda", [
    "Why Playwright?",
    "Setup & Installation",
    "Pytest Basics",
    "Playwright Test Structure",
    "Locators & Assertions",
    "Handling Windows & Dynamic Elements",
    "Record & Playback (Codegen)",
    "Q&A"
])

# Slide 3: Playwright Key Points
add_slide(bullet_slide_layout, "Playwright Key Points", [
    "What is Playwright?",
    "  Open-source E2E automation tool by Microsoft",
    "  Supports Chromium, Firefox, WebKit, Chrome/Edge",
    "Key Features & Benefits",
    "  Auto-Waiting: Waits for elements to be ready",
    "  Fast Execution: Parallel test runs",
    "  Cross-Browser Testing: Single script for all browsers",
    "  Codegen: Records actions to generate code",
    "  Network Interception: Mock/stub API calls",
    "  Advanced Scenarios: Popups, iframes, dynamic content",
    "  Bundled Browsers: No manual setup needed",
    "Why Choose Playwright?",
    "  Simplifies scalable, reliable web automation"
])

# Slide 4: Setup Instructions
add_slide(bullet_slide_layout, "Setup Instructions", [
    "Ensure Python 3.7+ is installed",
    "Install pytest:",
    "  pip install pytest",
    "Install pytest-playwright:",
    "  pip install pytest-playwright",
    "Install Playwright browsers:",
    "  playwright install"
])

# Slide 5: Pytest Basics
add_slide(bullet_slide_layout, "Pytest Basics", [
    "What is Pytest?",
    "  Python framework for unit, integration, E2E tests",
    "  Simple, readable, scalable with Playwright",
    "Syntax",
    "  Test files: test_*.py",
    "  Test functions: def test_example():",
    "Running Tests",
    "  Single file: pytest test_math.py",
    "  All tests: pytest",
    "  Specific test: pytest test_math.py::test_browser",
    "  Show print: pytest -s test_math.py"
])

# Slide 6: Pytest Fixture Concept
add_slide(bullet_slide_layout, "Pytest Fixture Concept", [
    "Reusable setup/teardown logic",
    "Declared with @pytest.fixture",
    "Example:",
    "  import pytest",
    "  @pytest.fixture",
    "  def open_browser():",
    "      print('Opening browser')",
    "      return 'browser_instance'",
    "  def test_browser(open_browser):",
    "      assert open_browser == 'browser_instance'",
    "Global Fixtures: Define in conftest.py"
])

# Slide 7: Playwright with Pytest
add_slide(bullet_slide_layout, "Playwright with Pytest", [
    "Basic Test with Built-in Fixtures",
    "  def test_playwrightbasics(playwright):",
    "      browser = playwright.chromium.launch()",
    "      context = browser.new_context()",
    "      page = context.new_page()",
    "      page.goto('{Url}')",
    "Headed Mode",
    "  Use: playwright.chromium.launch(headless=False)",
    "  CLI: pytest test_file.py --headed",
    "Enable Keyword Suggestions",
    "  from playwright.sync_api import Page",
    "  def test_example(page: Page):"
])

# Slide 8: Locators
add_slide(bullet_slide_layout, "Locators", [
    "get_by_label: For <label> tags",
    "  page.get_by_label('Username:').fill('rahulshettyacademy')",
    "get_by_role: For buttons, links, checkboxes",
    "  page.get_by_role('button', name='Sign In').click()",
    "get_by_text: Matches visible text",
    "  page.get_by_text('Incorrect username/password.').to_be_visible()",
    "get_by_placeholder: For placeholder attributes",
    "  expect(page.get_by_placeholder('Hide/Show Example')).to_be_visible()",
    "CSS Selectors: Fallback for complex cases",
    "  page.locator('#terms').check()"
])

# Slide 9: Assertions
add_slide(bullet_slide_layout, "Assertions", [
    "Auto-retrying to avoid flakiness",
    "  from playwright.sync_api import expect",
    "  expect(locator).to_be_checked()",
    "  expect(locator).to_be_visible()",
    "  expect(child_page.locator('.red')).to_contain_text('mentor@rahulshettyacademy.com')"
])

# Slide 10: Advanced Features
add_slide(bullet_slide_layout, "Advanced Features", [
    "Filtering Elements Dynamically",
    "  SamsungElem = page.locator('app-card').filter(has_text='Samsung Note 8')",
    "  SamsungElem.get_by_role('button').click()",
    "Codegen for Record and Playback",
    "  Run: playwright codegen {url}",
    "  Steps:",
    "    Open Chromium & Playwright Inspector",
    "    Perform actions manually",
    "    Stop recording; copy code from Inspector",
    "    Review and integrate into tests"
])

# Slide 11: Key Takeaways & Q&A
add_slide(bullet_slide_layout, "Key Takeaways & Q&A", [
    "Key Takeaways",
    "  Pytest + Playwright: Fast, scalable E2E testing",
    "  Fixtures streamline setup and reusability",
    "  Assertions reduce flakiness with auto-retry",
    "  Codegen for quick script generation",
    "Q&A",
    "  Questions? Ask away!"
])

# Save the presentation
prs.save("playwright_kt_presentation.pptx")